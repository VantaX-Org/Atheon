#!/usr/bin/env python3
"""Atheon (by Vantax) — premium 3-slide customer deck.

CINEMATIC DARK "flow river" hero. Each slide gets a full-bleed 2560x1440
PNG rendered with PIL: deep-navy field, a glowing three-layer bezier river
(halo / body / core) with color flowing blue -> teal -> green, and glowing
particle dots. Content sits on glassy translucent tiles (fnode look).
"""

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- brand tokens (DARK theme, from src/x/tokens.css) ----
BG_NAVY = (0x0a, 0x0d, 0x17)          # deep navy background
PANEL_A = (0x10, 0x15, 0x2a)          # panel low
PANEL_B = (0x14, 0x1b, 0x32)          # panel high
INK     = RGBColor(0xee, 0xf1, 0xfb)  # headlines
BODY     = RGBColor(0xa6, 0xae, 0xcb) # body text
MUT     = RGBColor(0x8b, 0x93, 0xb3)  # muted / kicker

GATE = (0x86, 0xa3, 0xff)   # accent blue
REC  = (0x2f, 0xe3, 0xa0)   # recovered green
FEE  = (0x3f, 0xb3, 0xac)   # fee teal
LEAK = (0xff, 0xc2, 0x4d)   # leak amber
REV  = (0xff, 0x64, 0x80)   # reversed pink

def _rgb(t):
    return RGBColor(t[0], t[1], t[2])

INK_GATE = _rgb(GATE)
INK_REC  = _rgb(REC)
INK_FEE  = _rgb(FEE)
INK_LEAK = _rgb(LEAK)
INK_REV  = _rgb(REV)

HEAD_FONT = "Schibsted Grotesk"
BODY_FONT = "IBM Plex Sans"
MONO_FONT = "Space Mono"

SW, SH = Inches(13.333), Inches(7.5)
PXW, PXH = 2560, 1440
# Live-system screenshots ship alongside this script; generated river
# backgrounds are transient, so they go to a temp dir.
import tempfile
ASSETS = Path(__file__).resolve().parent / "deck-assets"
SCRATCH = Path(tempfile.gettempdir()) / "atheon-deck-bg"
SCRATCH.mkdir(parents=True, exist_ok=True)


# ======================================================================
# PIL river background renderer
# ======================================================================
def lerp(a, b, t):
    return tuple(round(a[i] + (b[i] - a[i]) * t) for i in range(3))


def flow_color(t):
    """Blue -> teal -> green along the stream (t in 0..1)."""
    if t < 0.55:
        return lerp(GATE, FEE, t / 0.55)
    return lerp(FEE, REC, (t - 0.55) / 0.45)


def bezier(p0, p1, p2, p3, n):
    pts = []
    for i in range(n + 1):
        t = i / n
        mt = 1 - t
        x = (mt**3 * p0[0] + 3 * mt**2 * t * p1[0]
             + 3 * mt * t**2 * p2[0] + t**3 * p3[0])
        y = (mt**3 * p0[1] + 3 * mt**2 * t * p1[1]
             + 3 * mt * t**2 * p2[1] + t**3 * p3[1])
        pts.append((x, y))
    return pts


def stroke(layer, pts, width, alpha, blur):
    """Draw a color-flowing stroke on its own RGBA layer, blur it, return it."""
    d = ImageDraw.Draw(layer)
    n = len(pts) - 1
    r = width / 2
    for i in range(n):
        col = flow_color(i / n) + (alpha,)
        x0, y0 = pts[i]
        x1, y1 = pts[i + 1]
        d.line([(x0, y0), (x1, y1)], fill=col, width=int(width))
        # round the joints so thick strokes read smooth
        d.ellipse([x0 - r, y0 - r, x0 + r, y0 + r], fill=col)
    d.ellipse([pts[-1][0] - r, pts[-1][1] - r,
               pts[-1][0] + r, pts[-1][1] + r],
              fill=flow_color(1.0) + (alpha,))
    if blur:
        layer = layer.filter(ImageFilter.GaussianBlur(blur))
    return layer


def render_bg(path, curves, particles=44, vignette=0.10):
    """curves: list of (p0,p1,p2,p3, scale) control-point tuples in px."""
    base = Image.new("RGB", (PXW, PXH), BG_NAVY)

    # subtle radial vignette — lighter center
    if vignette:
        vg = Image.new("L", (PXW, PXH), 0)
        vd = ImageDraw.Draw(vg)
        cx, cy = PXW * 0.5, PXH * 0.46
        maxr = math.hypot(PXW, PXH) * 0.55
        # cheap radial: concentric ellipses
        steps = 60
        for s in range(steps, 0, -1):
            rr = maxr * s / steps
            v = int(255 * (1 - s / steps))
            vd.ellipse([cx - rr, cy - rr, cx + rr, cy + rr], fill=v)
        vg = vg.filter(ImageFilter.GaussianBlur(120))
        lighter = Image.new("RGB", (PXW, PXH),
                            lerp(BG_NAVY, (0x1a, 0x22, 0x40), vignette / 0.10))
        base = Image.composite(lighter, base, vg.point(lambda p: int(p * vignette)))

    for (p0, p1, p2, p3, scale) in curves:
        pts = bezier(p0, p1, p2, p3, 220)
        # three-layer glowing stroke, exactly like the app
        halo = stroke(Image.new("RGBA", (PXW, PXH), (0, 0, 0, 0)),
                      pts, 150 * scale, 60, blur=70)
        base = Image.alpha_composite(base.convert("RGBA"), halo).convert("RGB")
        body = stroke(Image.new("RGBA", (PXW, PXH), (0, 0, 0, 0)),
                      pts, 46 * scale, 150, blur=10)
        base = Image.alpha_composite(base.convert("RGBA"), body).convert("RGB")
        core = stroke(Image.new("RGBA", (PXW, PXH), (0, 0, 0, 0)),
                      pts, 8 * scale, 235, blur=1)
        base = Image.alpha_composite(base.convert("RGBA"), core).convert("RGB")

        # glowing particle dots along the curve
        glow = Image.new("RGBA", (PXW, PXH), (0, 0, 0, 0))
        gd = ImageDraw.Draw(glow)
        dot = Image.new("RGBA", (PXW, PXH), (0, 0, 0, 0))
        dd = ImageDraw.Draw(dot)
        for k in range(particles):
            t = (k + 0.5) / particles
            idx = int(t * (len(pts) - 1))
            x, y = pts[idx]
            col = flow_color(t)
            # size pulses a little along the path
            r = 6 + 6 * abs(math.sin(t * 9.0))
            gr = r * 3.2
            gd.ellipse([x - gr, y - gr, x + gr, y + gr], fill=col + (110,))
            dd.ellipse([x - r, y - r, x + r, y + r], fill=col + (255,))
        glow = glow.filter(ImageFilter.GaussianBlur(16))
        base = Image.alpha_composite(base.convert("RGBA"), glow).convert("RGB")
        base = Image.alpha_composite(base.convert("RGBA"), dot).convert("RGB")

    base.save(path)
    return path


# ======================================================================
# pptx primitives
# ======================================================================
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


def no_shadow(shp):
    shp.shadow.inherit = False


def full_bg(slide, png):
    pic = slide.shapes.add_picture(str(png), 0, 0, SW, SH)
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)
    return pic


def set_fill_alpha(shp, alpha):
    """Set solidFill transparency via <a:alpha> (python-pptx has no API)."""
    spPr = shp._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    srgb = sf.find(qn('a:srgbClr'))
    for old in srgb.findall(qn('a:alpha')):
        srgb.remove(old)
    a = srgb.makeelement(qn('a:alpha'), {'val': str(alpha)})
    srgb.append(a)


def set_line_alpha(shp, alpha):
    ln = shp._element.spPr.find(qn('a:ln'))
    sf = ln.find(qn('a:solidFill'))
    srgb = sf.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(alpha)})
    srgb.append(a)


def soft_shadow(shp, blur=120000, dist=55000, alpha=52000):
    spPr = shp._element.spPr
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
                        {'blurRad': str(blur), 'dist': str(dist),
                         'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '000000'})
    clr.append(clr.makeelement(qn('a:alpha'), {'val': str(alpha)}))
    sh.append(clr); el.append(sh); spPr.append(el)


def glass(slide, x, y, w, h, radius=0.09, panel=PANEL_A, alpha=62000,
          border=(0x94, 0xa3, 0xd6), border_alpha=22000, shadow=True):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        c.adjustments[0] = radius
    except Exception:
        pass
    c.fill.solid(); c.fill.fore_color.rgb = _rgb(panel)
    set_fill_alpha(c, alpha)
    c.line.color.rgb = _rgb(border); c.line.width = Pt(1)
    set_line_alpha(c, border_alpha)
    no_shadow(c)
    if shadow:
        soft_shadow(c)
    return c


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        spacing=1.0, para_gap=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if para_gap:
            p.space_after = Pt(para_gap)
        for (t, font, size, color, bold, track) in para:
            r = p.add_run(); r.text = t
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold
            if track:
                r.font._rPr.set('spc', str(int(track * 100)))
    return tb


def kicker(slide, x, y, label, color=MUT):
    txt(slide, x, y, Inches(9), Inches(0.3),
        [[(label.upper(), MONO_FONT, 11, color, True, 2.4)]])


def chip(slide, x, y, color):
    """Small colored node dot, top-left of a tile."""
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.16), Inches(0.16))
    d.fill.solid(); d.fill.fore_color.rgb = _rgb(color)
    d.line.fill.background(); no_shadow(d)
    soft_shadow(d, blur=60000, dist=0, alpha=60000)
    return d


def screenshot_slide(prs, png, bg, kick, title, lede, points, accent=GATE):
    """A real-system screenshot on the right, an explanatory caption column
    on the left. points: list of (heading, body) bullet tiles."""
    s = prs.slides.add_slide(BLANK)
    full_bg(s, bg)

    # ---- screenshot on the right, in a glowing glass frame ----
    img_w = Inches(8.15)
    img_h = Emu(int(int(img_w) * 2000 / 3200))   # captured 16:10
    img_x = Emu(int(SW) - int(img_w) - int(Inches(0.55)))
    img_y = Emu(int((int(SH) - int(img_h)) / 2))
    # frame glow behind the shot
    fr = glass(s, Emu(int(img_x) - Inches(0.12)), Emu(int(img_y) - Inches(0.12)),
               Emu(int(img_w) + Inches(0.24)), Emu(int(img_h) + Inches(0.24)),
               radius=0.03, panel=PANEL_B, alpha=70000,
               border=accent, border_alpha=40000)
    soft_shadow(fr, blur=180000, dist=70000, alpha=60000)
    pic = s.shapes.add_picture(str(png), img_x, img_y, img_w, img_h)
    pic.line.color.rgb = _rgb(accent); pic.line.width = Pt(1)
    set_line_alpha(pic, 30000)

    # ---- caption column on the left ----
    cx = Inches(0.62)
    cw = Emu(int(img_x) - int(cx) - int(Inches(0.35)))
    kicker(s, cx, Inches(0.72), kick, color=_rgb(accent))
    txt(s, cx, Inches(1.04), cw, Inches(1.3),
        [[(title, HEAD_FONT, 25, INK, True, 0)]], spacing=1.02)
    txt(s, cx, Inches(2.28), cw, Inches(1.1),
        [[(lede, BODY_FONT, 12.5, MUT, False, 0)]], spacing=1.26)

    py = Inches(3.62)
    for (head, body) in points:
        chip(s, cx, Emu(int(py) + Inches(0.05)), accent)
        txt(s, Emu(int(cx) + Inches(0.30)), py, cw, Inches(0.34),
            [[(head, HEAD_FONT, 13, INK, True, 0)]])
        txt(s, Emu(int(cx) + Inches(0.30)), Emu(int(py) + Inches(0.32)),
            Emu(int(cw) - Inches(0.30)), Inches(0.7),
            [[(body, BODY_FONT, 11, BODY, False, 0)]], spacing=1.18)
        py = Emu(int(py) + Inches(1.02))
    return s


# ======================================================================
# render backgrounds
# ======================================================================
BG1 = SCRATCH / "bg1.png"
BG2 = SCRATCH / "bg2.png"
BG3 = SCRATCH / "bg3.png"

# Slide 1: broad flowing river through lower-middle where the tiles sit
render_bg(BG1, [
    ((-120, 820), (700, 560), (1750, 1120), (2700, 760), 1.25),
    ((-120, 980), (820, 1180), (1650, 620), (2700, 980), 0.7),
], particles=50)

# Slide 2: single high sweeping stream across the top
render_bg(BG2, [
    ((-120, 300), (760, 120), (1780, 520), (2700, 240), 1.0),
], particles=40)

# Slide 3: two streams converging low
render_bg(BG3, [
    ((-120, 900), (760, 700), (1720, 1080), (2700, 820), 1.15),
    ((-120, 1080), (820, 1240), (1650, 760), (2700, 1060), 0.65),
], particles=46)

# Screenshot slides: a quieter single stream low on the left, so it lives
# under the caption column and doesn't fight the screenshot on the right.
BGS = SCRATCH / "bgs.png"
render_bg(BGS, [
    ((-160, 1120), (520, 900), (1200, 1320), (1900, 980), 0.8),
], particles=26, vignette=0.08)


# ======================================================================
# SLIDE 1 — How it works
# ======================================================================
s = prs.slides.add_slide(BLANK)
full_bg(s, BG1)

kicker(s, Inches(0.7), Inches(0.55), "Atheon  ·  by Vantax", color=INK_GATE)
txt(s, Inches(0.7), Inches(0.85), Inches(12.0), Inches(1.1),
    [[("A flowing river of value — from signal to the boardroom",
       HEAD_FONT, 30, INK, True, 0)]])
txt(s, Inches(0.7), Inches(2.05), Inches(11.4), Inches(0.6),
    [[("External reads and your Catalysts flow into a value river that surfaces "
       "leakage across the operational value chain — then rolls up into "
       "management and executive views.", BODY_FONT, 14, MUT, False, 0)]],
    spacing=1.25)

nodes = [
    ("EXTERNAL", "Signals + Catalysts",
     "Market, regulatory & supplier reads meet your Catalysts", GATE),
    ("THE RIVER", "Value river",
     "Detects leakage priced across every value-chain stage", FEE),
    ("ROLL-UP", "Management view",
     "Where value leaks, what's recoverable, what's decided", LEAK),
    ("BOARDROOM", "Executive view",
     "One honest number, drill-through to the transaction", REC),
]
n = len(nodes)
cw, gap = Inches(2.78), Inches(0.28)
total = cw * n + gap * (n - 1)
x0 = int((SW - total) / 2)
cy, ch = Inches(2.95), Inches(2.02)
for i, (tag, title, desc, accent) in enumerate(nodes):
    x = Emu(x0 + i * int(cw + gap))
    glass(s, x, cy, cw, ch, radius=0.08)
    chip(s, Emu(int(x) + Inches(0.30)), Emu(int(cy) + Inches(0.32)), accent)
    txt(s, Emu(int(x) + Inches(0.56)), Emu(int(cy) + Inches(0.27)),
        Emu(int(cw) - Inches(0.7)), Inches(0.3),
        [[(tag, MONO_FONT, 9.5, MUT, True, 1.8)]])
    txt(s, Emu(int(x) + Inches(0.30)), Emu(int(cy) + Inches(0.72)),
        Emu(int(cw) - Inches(0.56)), Inches(0.6),
        [[(title, HEAD_FONT, 17, INK, True, 0)]], spacing=1.02)
    txt(s, Emu(int(x) + Inches(0.30)), Emu(int(cy) + Inches(1.42)),
        Emu(int(cw) - Inches(0.58)), Inches(0.55),
        [[(desc, BODY_FONT, 11, BODY, False, 0)]], spacing=1.16)

# Jeff strip
jy = Inches(5.45)
glass(s, Inches(0.7), jy, Inches(11.93), Inches(1.32), radius=0.10,
      panel=PANEL_B, alpha=66000, border=GATE, border_alpha=34000)
d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05),
                       Emu(int(jy) + Inches(0.36)), Inches(0.6), Inches(0.6))
d.fill.solid(); d.fill.fore_color.rgb = _rgb(GATE)
d.line.fill.background(); no_shadow(d)
soft_shadow(d, blur=90000, dist=0, alpha=45000)
txt(s, Inches(1.05), Emu(int(jy) + Inches(0.32)), Inches(0.6), Inches(0.68),
    [[("✦", BODY_FONT, 24, _rgb(BG_NAVY), True, 0)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.95), Emu(int(jy) + Inches(0.26)), Inches(10.4), Inches(0.4),
    [[("Meet Jeff — the grounded AI assistant", HEAD_FONT, 16, INK, True, 0)]])
txt(s, Inches(1.95), Emu(int(jy) + Inches(0.64)), Inches(10.5), Inches(0.6),
    [[("Jeff explains every figure in plain language and answers only from your "
       "tenant's real data — he never invents numbers. Ask any node, get the "
       "booked field behind it.", BODY_FONT, 12, BODY, False, 0)]], spacing=1.2)


# ======================================================================
# SCREENSHOT WALKTHROUGH — the real, live system
# ======================================================================
screenshot_slide(
    prs, ASSETS / "sys-login.png", BGS,
    "The live system  ·  first screen",
    "The brand is the river, from the login on",
    "No stock hero, no placeholder. The sign-in screen already tells the story — "
    "Connect, Detect, Fix, Recover, Report — as one flowing current. The same "
    "visual language runs through every surface behind it.",
    [
        ("The money is already in your ERP",
         "The promise is on the door: Atheon recovers value you have already earned but are quietly leaking."),
        ("One design language",
         "The glowing river you sign in against is the exact motif that carries through the console, the ledger and Jeff."),
    ],
    accent=GATE,
)

screenshot_slide(
    prs, ASSETS / "sys-cfo.png", BGS,
    "The live system  ·  the console",
    "One console — the whole business, as a river",
    "The recovered number leads. Beneath it, your value chain flows left-to-right — "
    "each stage priced from booked fields — then crosses the boundary into leakage "
    "detected, decisions awaiting a signature, and money recovered.",
    [
        ("Every figure is real",
         "R6.96m recovered, R1.45m leakage this assessment, R4.6m awaiting sign-off — each traces to a booked API field, never a mock."),
        ("Signal to boardroom in one view",
         "Live external signals scroll above the river; the value chain, the recovery flow and the decision queue all sit on one screen."),
        ("Health, anomalies, risk at a glance",
         "The pulse strip carries health, red metrics, open anomalies and risk alerts since the last period."),
    ],
    accent=REC,
)

screenshot_slide(
    prs, ASSETS / "sys-coo.png", BGS,
    "The live system  ·  role lenses",
    "Every seat gets its own river",
    "Switch the seat and the whole console re-lenses. The COO's chain reads "
    "Source & Plan · Inbound & Stores · Production & Delivery · People & Shifts · "
    "Ship & Bill — the CFO's and CPO's name the same flow in their own terms.",
    [
        ("The chain is relabelled per role",
         "Same booked data, named the way each executive thinks — operations, procurement or finance — not one generic layout for all."),
        ("The lede changes with the seat",
         "“Where operations leak value and which catalysts are running on it” for the COO; supplier-side leakage first for the CPO."),
        ("Approval rights follow the role",
         "Who can sign a decision off is set by role and lens — the API stays the enforcement point."),
    ],
    accent=FEE,
)

screenshot_slide(
    prs, ASSETS / "sys-leak.png", BGS,
    "The live system  ·  drill-through",
    "From any stage, straight to the finding",
    "Click a stage on the river and it opens — the priced leakage, the health "
    "trend, the downstream stages it feeds, and where the number came from. "
    "Provenance is stated, not implied.",
    [
        ("Impact downstream, made explicit",
         "This node feeds Cash & ledger (R102k, 2 findings) and Tax & filings (R235k, 3 findings) — the blast radius, not just the local number."),
        ("Provenance on every figure",
         "“Priced from the assessment findings summary — an estimate until recovered and booked.” The system says how sure it is."),
        ("Two ways deeper",
         "Open the findings for the transactions, or ask Jeff to explain it in plain language — both one click away."),
    ],
    accent=LEAK,
)

screenshot_slide(
    prs, ASSETS / "sys-jeff.png", BGS,
    "The live system  ·  grounded AI",
    "Jeff — grounded answers, now with a voice",
    "Ask any figure in plain language. Jeff answers only from your tenant's real "
    "data — the ROI summary, active risks, recent catalyst runs — and names the "
    "sources he used. He never invents a number.",
    [
        ("Real findings, cited",
         "Inventory Reconciliation #6 (ZAR 214k, 16 exceptions), GR/IR #6 (ZAR 422.5k), AP Invoice Validation #6 (ZAR 306k) — pulled from booked runs."),
        ("Talk to Jeff, or let him talk back",
         "A microphone dictates your question and Jeff can read his answer aloud — hands-free in the boardroom."),
        ("Shows his working",
         "Every reply footers the data it used and the edge latency — “Atheon Edge · 13191 ms” — so the answer is auditable."),
    ],
    accent=GATE,
)


# ======================================================================
# SLIDE 2 — Features
# ======================================================================
s = prs.slides.add_slide(BLANK)
full_bg(s, BG2)

kicker(s, Inches(0.7), Inches(0.52), "What you get", color=INK_GATE)
txt(s, Inches(0.7), Inches(0.82), Inches(11.9), Inches(0.8),
    [[("The whole recovery loop, on one river", HEAD_FONT, 31, INK, True, 0)]])

feats = [
    ("Live Recovery Console",
     "The flow river on every surface — signals stream in, leakage surfaces, decisions settle.", GATE),
    ("External signal radar",
     "Market, regulatory and supplier reads, each with its own computed business impact.", GATE),
    ("Leakage detection",
     "Priced across the value chain, drill-down from any stage straight to the transactions.", LEAK),
    ("Decisions gate",
     "Every call passes a gate and lands a sealed, tamper-evident audit-chain receipt.", REV),
    ("Recovery ledger",
     "Money recovered, logged with reported ROI — traceable line by line.", REC),
    ("C-suite role lenses",
     "CEO, CFO, CPO and more — each seat sees its own river, tuned to its stakes.", FEE),
    ("Jeff AI drill-through",
     "Ask any figure; Jeff answers from booked data and cites the field or receipt.", GATE),
    ("One folded console",
     "Operations, assurance and admin folded into a single interface — no tab sprawl.", (0x8b, 0x93, 0xb3)),
]
cols = 4
cw = Inches(2.86); ch = Inches(2.30)
gx = Inches(0.20); gy = Inches(0.28)
grid_w = cw * cols + gx * (cols - 1)
x0 = int((SW - grid_w) / 2)
y0 = Inches(1.95)
for idx, (title, desc, accent) in enumerate(feats):
    r, cc = divmod(idx, cols)
    x = Emu(x0 + cc * int(cw + gx))
    y = Emu(int(y0) + r * int(ch + gy))
    glass(s, x, y, cw, ch, radius=0.08)
    chip(s, Emu(int(x) + Inches(0.32)), Emu(int(y) + Inches(0.34)), accent)
    txt(s, Emu(int(x) + Inches(0.32)), Emu(int(y) + Inches(0.72)),
        Emu(int(cw) - Inches(0.6)), Inches(0.7),
        [[(title, HEAD_FONT, 15.5, INK, True, 0)]], spacing=1.02)
    txt(s, Emu(int(x) + Inches(0.32)), Emu(int(y) + Inches(1.28)),
        Emu(int(cw) - Inches(0.62)), Inches(0.9),
        [[(desc, BODY_FONT, 11, BODY, False, 0)]], spacing=1.18)


# ======================================================================
# SLIDE 3 — The value to the organization
# ======================================================================
s = prs.slides.add_slide(BLANK)
full_bg(s, BG3)

kicker(s, Inches(0.7), Inches(0.52), "Why it matters", color=INK_GATE)
txt(s, Inches(0.7), Inches(0.82), Inches(11.9), Inches(0.8),
    [[("Value to the organization", HEAD_FONT, 31, INK, True, 0)]])
txt(s, Inches(0.7), Inches(1.62), Inches(11.4), Inches(0.5),
    [[("Money recovered, one honest source of truth, and faster decisions — "
       "every figure traces to a booked field or a sealed receipt.",
       BODY_FONT, 14, MUT, False, 0)]], spacing=1.2)

vals = [
    ("Money recovered",
     "Value that was quietly leaking, caught and booked back into the P&L.", REC),
    ("One honest truth",
     "No invented figures — everything traces to a booked field or sealed receipt.", REC),
    ("Faster decisions",
     "Executives decide with drill-down to the underlying transaction, in the same view.", GATE),
    ("Risk seen early",
     "External market, regulatory and supplier risk surfaced before it hits the P&L.", LEAK),
    ("Clarity per seat",
     "Role-relevant lenses give every C-suite seat its own river and its own stakes.", FEE),
]
n = len(vals)
cw = Inches(2.28); gap = Inches(0.18)
total = cw * n + gap * (n - 1)
x0 = int((SW - total) / 2)
cy = Inches(2.62); ch = Inches(3.0)
for i, (title, desc, accent) in enumerate(vals):
    x = Emu(x0 + i * int(cw + gap))
    glass(s, x, cy, cw, ch, radius=0.07)
    chip(s, Emu(int(x) + Inches(0.28)), Emu(int(cy) + Inches(0.34)), accent)
    txt(s, Emu(int(x) + Inches(0.28)), Emu(int(cy) + Inches(0.68)),
        Emu(int(cw) - Inches(0.52)), Inches(0.8),
        [[(title, HEAD_FONT, 17, INK, True, 0)]], spacing=1.02)
    txt(s, Emu(int(x) + Inches(0.28)), Emu(int(cy) + Inches(1.42)),
        Emu(int(cw) - Inches(0.52)), Inches(1.5),
        [[(desc, BODY_FONT, 11.5, BODY, False, 0)]], spacing=1.24)

txt(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.4),
    [[("Every claim above is a capability of the platform — Atheon reports only "
       "what your data proves.", MONO_FONT, 10, MUT, False, 0.6)]])


# ---------- save ----------
out = Path("/Users/reshigan/Atheon/docs/valuation/Atheon-Customer-Deck.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"OK {out}  slides={len(prs.slides._sldIdLst)}")
